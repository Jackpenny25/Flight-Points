#!/usr/bin/env python3
"""Generate a PowerPoint from site data (attendance, cadets).

Usage:
  python scripts/generate_presentation.py --data-dir data --output outputs/cadets_presentation.pptx

Features:
- Reads JSON files in `data/` (attendance.json, cadets.json)
- Creates title slide, attendance summary, top-cadets table, and a points chart
- Outputs a PPTX file and temporary chart images (in outputs/)
"""
import argparse
import json
import os
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


def ensure_dir(path: Path):
    path.mkdir(parents=True, exist_ok=True)


def load_json(path: Path):
    if not path.exists():
        return None
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def attendance_summary(attendance_data):
    # attendance_data expected as list of events each containing attendees list or mapping
    if not attendance_data:
        return {"events": 0, "last_attendance": 0, "avg_attendance": 0}
    try:
        events = attendance_data
        counts = []
        for ev in events:
            if isinstance(ev, dict):
                # try common shapes
                if "attendees" in ev and isinstance(ev["attendees"], list):
                    counts.append(len(ev["attendees"]))
                elif "present" in ev and isinstance(ev["present"], list):
                    counts.append(len(ev["present"]))
                elif "count" in ev and isinstance(ev["count"], int):
                    counts.append(ev["count"])                
                else:
                    # fallback: number of keys
                    counts.append(len(ev))
            elif isinstance(ev, list):
                counts.append(len(ev))
        events_count = len(counts)
        last = counts[-1] if counts else 0
        avg = sum(counts) / len(counts) if counts else 0
        return {"events": events_count, "last_attendance": last, "avg_attendance": round(avg, 1)}
    except Exception:
        return {"events": 0, "last_attendance": 0, "avg_attendance": 0}


def top_cadets_table(cadets_data, top_n=10):
    # cadets_data is expected to be list of cadet dicts with 'name' and 'points' fields
    if not cadets_data:
        return []
    df = pd.json_normalize(cadets_data)
    if "points" not in df.columns:
        # try to infer points column
        possible = [c for c in df.columns if "point" in c.lower()]
        pts = possible[0] if possible else None
    else:
        pts = "points"
    name_col = None
    for c in df.columns:
        if c.lower() in ("name", "fullname", "full_name"):
            name_col = c
            break
    if name_col is None and "id" in df.columns:
        name_col = "id"
    if pts is None:
        df["points"] = 0
        pts = "points"
    df = df.fillna({pts: 0})
    df_sorted = df.sort_values(by=pts, ascending=False)
    rows = []
    for _, r in df_sorted.head(top_n).iterrows():
        name = r.get(name_col, "Unknown")
        points = int(r.get(pts, 0))
        rows.append((str(name), points))
    return rows


def make_points_chart(top_rows, out_png: Path):
    names = [r[0] for r in top_rows]
    points = [r[1] for r in top_rows]
    plt.figure(figsize=(10, 4))
    bars = plt.barh(names[::-1], points[::-1], color="#1f77b4")
    plt.xlabel("Points")
    plt.title("Top Cadets by Points")
    for bar in bars:
        w = bar.get_width()
        plt.text(w + max(points) * 0.01, bar.get_y() + bar.get_height() / 2, f"{int(w)}", va="center")
    plt.tight_layout()
    plt.savefig(out_png, dpi=150)
    plt.close()


def add_title_slide(prs: Presentation, title: str, subtitle: str = None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_el = slide.shapes.title
    subtitle_el = slide.placeholders[1]
    title_el.text = title
    if subtitle:
        subtitle_el.text = subtitle
    return slide


def add_text_slide(prs: Presentation, title: str, paragraphs: list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, p in enumerate(paragraphs):
        if i == 0:
            body.text = p
        else:
            p_el = body.add_paragraph()
            p_el.text = p
            p_el.level = 0
    return slide


def add_image_slide(prs: Presentation, title: str, img_path: Path, width_inches=9):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.5)
    slide.shapes.add_picture(str(img_path), left, top, width=Inches(width_inches))
    return slide


def add_table_slide(prs: Presentation, title: str, rows: list, col_names=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if not rows:
        return slide
    cols = len(rows[0])
    rows_count = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 + 0.25 * rows_count)
    table_shape = slide.shapes.add_table(rows_count, cols, left, top, width, height)
    table = table_shape.table
    # header
    headers = col_names or ["Name", "Points"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            table.cell(r_idx, c_idx).text = str(val)
    return slide


def build_presentation(data_dir: Path, output: Path, title_prefix: str = "Cadet Update"):
    ensure_dir(output.parent)
    attendance = load_json(data_dir / "attendance.json")
    cadets = load_json(data_dir / "cadets.json")

    summary = attendance_summary(attendance)
    top_rows = top_cadets_table(cadets, top_n=10)

    prs = Presentation()
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    add_title_slide(prs, f"{title_prefix} — {datetime.now().strftime('%d %b %Y')}", subtitle=f"Generated: {timestamp}")

    paragraphs = [
        f"Events considered: {summary['events']}",
        f"Last meeting attendance: {summary['last_attendance']}",
        f"Average attendance: {summary['avg_attendance']}"
    ]
    add_text_slide(prs, "Attendance Summary", paragraphs)

    if top_rows:
        chart_png = output.parent / "top_points.png"
        make_points_chart(top_rows, chart_png)
        add_image_slide(prs, "Top Cadets — Points", chart_png)
        add_table_slide(prs, "Top Cadets Table", top_rows, col_names=["Name", "Points"])

    prs.save(str(output))
    return output


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--data-dir", default="data", help="Directory with attendance.json and cadets.json")
    parser.add_argument("--output", default="outputs/cadets_presentation.pptx", help="Output PPTX file")
    parser.add_argument("--title", default="Cadet Update", help="Title prefix for the slides")
    args = parser.parse_args()

    data_dir = Path(args.data_dir)
    output = Path(args.output)
    out_file = build_presentation(data_dir, output, title_prefix=args.title)
    print(f"Generated: {out_file}")


if __name__ == "__main__":
    main()
