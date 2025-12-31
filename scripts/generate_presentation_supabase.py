#!/usr/bin/env python3
"""Generate a PowerPoint by fetching data from Supabase Functions.

This script will attempt to call the project's Supabase Functions endpoints (deployed function "server").
It prefers using a service role key (SUPABASE_SERVICE_ROLE_KEY) for full access. If only public data is available,
it will fall back to public endpoints.

Environment variables (recommended):
- SUPABASE_PROJECT_ID (optional) - defaults to the project id in utils/supabase/info.tsx
- SUPABASE_FUNCTION_NAME (optional) - defaults to "server"
- SUPABASE_SERVICE_ROLE_KEY (recommended) - service key for protected endpoints
- SUPABASE_ANON_KEY (optional) - anon key for public endpoints

Usage:
    python scripts/generate_presentation_supabase.py --output outputs/cadets_presentation.pptx
"""
from __future__ import annotations
import argparse
import json
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

import matplotlib.pyplot as plt
import pandas as pd
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from urllib.parse import quote


REPO_ROOT = Path(__file__).resolve().parent.parent


def read_project_id_from_repo() -> Optional[str]:
    info_file = REPO_ROOT / "utils" / "supabase" / "info.tsx"
    if not info_file.exists():
        return None
    text = info_file.read_text(encoding="utf-8")
    m = re.search(r"projectId\s*=\s*\"([a-z0-9]+)\"", text)
    if m:
        return m.group(1)
    return None


def build_functions_base(project_id: str, function_name: str = "server") -> str:
    return f"https://{project_id}.functions.supabase.co/{function_name}"


def request_json(url: str, api_key: Optional[str] = None) -> Optional[Dict[str, Any]]:
    headers = {"Accept": "application/json"}
    if api_key:
        headers.update({"apikey": api_key, "Authorization": f"Bearer {api_key}"})
    try:
        r = requests.get(url, headers=headers, timeout=15)
        r.raise_for_status()
        return r.json()
    except Exception as e:
        print(f"Warning: request to {url} failed: {e}")
        return None


def try_endpoints(base: str, paths: List[str], api_key: Optional[str]) -> Optional[Dict[str, Any]]:
    for p in paths:
        url = base.rstrip("/") + "/" + p.lstrip("/")
        j = request_json(url, api_key)
        if j is not None:
            return j
    return None


def fetch_cadets(base: str, api_key: Optional[str]) -> List[Dict[str, Any]]:
    # Try public cadets endpoints then common variants
    paths = ["public/cadets", "cadets", "make-server-73a3871f/public/cadets", "make-server-73a3871f/cadets"]
    j = try_endpoints(base, paths, api_key)
    if not j:
        return []
    # functions return { cadets: [...] } or raw array
    if isinstance(j, dict) and "cadets" in j:
        return j["cadets"] or []
    if isinstance(j, list):
        return j
    return []


def fetch_leaderboards(base: str, api_key: Optional[str]) -> Dict[str, Any]:
    paths = ["make-server-73a3871f/leaderboards", "leaderboards"]
    j = try_endpoints(base, paths, api_key)
    if not j:
        return {}
    return j


def fetch_kv_table_postgrest(project_id: str, table: str, api_key: Optional[str], prefix: Optional[str] = None) -> List[Dict[str, Any]]:
    """Fetch rows from the KV table via Supabase PostgREST using the service role key.
    Returns list of {'key':..., 'value': ...} rows.
    """
    if not api_key:
        return []
    base = f"https://{project_id}.supabase.co/rest/v1/{table}"
    if prefix:
        # build a like filter for keys that start with prefix
        # PostgREST expects percent-encoded % as %25
        filter_expr = f"key=like.%25{quote(prefix)}%25"
    else:
        filter_expr = ""
    url = base + "?select=key,value"
    if filter_expr:
        url = url + "&" + filter_expr
    headers = {
        "apikey": api_key,
        "Authorization": f"Bearer {api_key}",
        "Accept": "application/json",
    }
    try:
        r = requests.get(url, headers=headers, timeout=20)
        r.raise_for_status()
        return r.json() or []
    except Exception as e:
        print(f"Warning: PostgREST request to {url} failed: {e}")
        return []


def attendance_summary_from_api(base: str, api_key: Optional[str]) -> Dict[str, Any]:
    # Try several attendance report endpoints
    paths = [
        "make-server-73a3871f/attendance/reports",
        "attendance/reports",
        "make-server-73a3871f/attendance",
        "attendance",
    ]
    j = try_endpoints(base, paths, api_key)
    if not j:
        return {"events": 0, "last_attendance": 0, "avg_attendance": 0}
    # If the function returns { summary, stats } or { attendance: [...] }
    if isinstance(j, dict):
        if "stats" in j:
            stats = j["stats"]
            return {"events": stats.get("totalRecords", 0), "last_attendance": stats.get("totalPresent", 0), "avg_attendance": stats.get("averageAttendanceRate", 0)}
        if "attendance" in j and isinstance(j["attendance"], list):
            counts = len(j["attendance"])
            last = len(j["attendance"][0].get("present", []) if j["attendance"] else [])
            return {"events": counts, "last_attendance": last, "avg_attendance": 0}
    return {"events": 0, "last_attendance": 0, "avg_attendance": 0}


def top_cadets_from_leaderboard(lb: Dict[str, Any], top_n: int = 10) -> List[tuple[str, int]]:
    rows: List[tuple[str, int]] = []
    if not lb:
        return rows
    cadets = lb.get("cadetLeaderboard") or lb.get("cadetLeaderboard")
    if isinstance(cadets, list):
        for item in cadets[:top_n]:
            if isinstance(item, dict):
                name = item.get("name") or item.get("cadetName") or item.get("id")
                pts = item.get("points") or item.get("points", 0)
                rows.append((str(name), int(pts or 0)))
    return rows


# Presentation helper functions (reused from local script)

def make_points_chart(top_rows, out_png: Path):
    names = [r[0] for r in top_rows]
    points = [r[1] for r in top_rows]
    if not names:
        return None
    plt.figure(figsize=(10, 4))
    bars = plt.barh(names[::-1], points[::-1], color="#1f77b4")
    plt.xlabel("Points")
    plt.title("Top Cadets by Points")
    maxp = max(points) if points else 1
    for bar in bars:
        w = bar.get_width()
        plt.text(w + maxp * 0.01, bar.get_y() + bar.get_height() / 2, f"{int(w)}", va="center")
    plt.tight_layout()
    plt.savefig(out_png, dpi=150)
    plt.close()


def add_title_slide(prs: Presentation, title: str, subtitle: str = None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_el = slide.shapes.title
    subtitle_el = slide.placeholders[1]
    title_el.text = title
    if subtitle:
        subtitle_el.text = subtitle
    return slide


def add_text_slide(prs: Presentation, title: str, paragraphs: list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, p in enumerate(paragraphs):
        if i == 0:
            body.text = p
        else:
            p_el = body.add_paragraph()
            p_el.text = p
            p_el.level = 0
    return slide


def add_image_slide(prs: Presentation, title: str, img_path: Path, width_inches=9):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.5)
    slide.shapes.add_picture(str(img_path), left, top, width=Inches(width_inches))
    return slide


def add_table_slide(prs: Presentation, title: str, rows: list, col_names=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if not rows:
        return slide
    cols = len(rows[0])
    rows_count = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 + 0.25 * rows_count)
    table_shape = slide.shapes.add_table(rows_count, cols, left, top, width, height)
    table = table_shape.table
    headers = col_names or ["Name", "Points"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            table.cell(r_idx, c_idx).text = str(val)
    return slide


def build_presentation_from_supabase(output: Path, project_id: str, function_name: str, api_key: Optional[str]):
    base = build_functions_base(project_id, function_name)

    cadets = []
    lb = {}
    summary = {"events": 0, "last_attendance": 0, "avg_attendance": 0}
    top_rows: List[tuple[str, int]] = []

    # If we have a service role key, try PostgREST KV fallback first (avoids 401 on Functions)
    if api_key:
        print("Attempting PostgREST KV fallback using service role key...")
        kv_rows = fetch_kv_table_postgrest(project_id, "kv_store_73a3871f", api_key, prefix="cadet:")
        if kv_rows:
            print(f"PostgREST: found {len(kv_rows)} KV rows. Parsing cadets...")
            for r in kv_rows:
                val = r.get("value")
                try:
                    parsed = json.loads(val)
                    cadets.append(parsed)
                except Exception:
                    cadets.append({"id": r.get("key"), "name": val})
        else:
            print("PostgREST: no KV cadet rows found or PostgREST request failed; will try Functions endpoints next.")

    # Try Functions endpoints if cadets not populated
    if not cadets:
        cadets = fetch_cadets(base, api_key)
        lb = fetch_leaderboards(base, api_key) if api_key else {}
        summary = attendance_summary_from_api(base, api_key)
        top_rows = top_cadets_from_leaderboard(lb, top_n=10) if lb else []
    else:
        # If cadets found via KV, attempt to compute leaderboard from cadets values if possible
        df = pd.json_normalize(cadets)
        if not df.empty:
            pts_col = None
            for c in df.columns:
                if "point" in c.lower():
                    pts_col = c
                    break
            name_col = None
            for c in df.columns:
                if c.lower() in ("name", "fullname", "full_name"):
                    name_col = c
                    break
            if pts_col and name_col:
                df = df.fillna({pts_col: 0})
                df_sorted = df.sort_values(by=pts_col, ascending=False)
                for _, r in df_sorted.head(10).iterrows():
                    top_rows.append((r.get(name_col, "Unknown"), int(r.get(pts_col, 0))))

    prs = Presentation()
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
    add_title_slide(prs, f"Cadet Update — {datetime.now().strftime('%d %b %Y')}", subtitle=f"Generated: {timestamp}")

    paragraphs = [
        f"Cadets fetched: {len(cadets)}",
        f"Events considered: {summary.get('events', 0)}",
        f"Last meeting attendance: {summary.get('last_attendance', 0)}",
        f"Average attendance: {summary.get('avg_attendance', 0)}",
    ]
    add_text_slide(prs, "Attendance Summary", paragraphs)
    # ensure output directory exists before saving (prevents FileNotFoundError)
    try:
        output.parent.mkdir(parents=True, exist_ok=True)
    except Exception:
        pass

    if not cadets:
        # fallback: try reading kv table directly using PostgREST and the service role key
        kv_table = "kv_store_73a3871f"
        rows = fetch_kv_table_postgrest(project_id, kv_table, api_key, prefix="cadet:")
        cadets = []
        for r in rows:
            val = r.get("value")
            try:
                parsed = json.loads(val)
                cadets.append(parsed)
            except Exception:
                cadets.append({"id": r.get("key"), "name": val})

    if top_rows:
        chart_png = output.parent / "top_points_supabase.png"
        make_points_chart(top_rows, chart_png)
        add_image_slide(prs, "Top Cadets — Points", chart_png)
        add_table_slide(prs, "Top Cadets Table", top_rows, col_names=["Name", "Points"])
    else:
        # If no leaderboard available, try to compute from cadets if they have points
        df = pd.json_normalize(cadets)
        if not df.empty:
            pts_col = None
            for c in df.columns:
                if "point" in c.lower():
                    pts_col = c
                    break
            name_col = None
            for c in df.columns:
                if c.lower() in ("name", "fullname", "full_name"):
                    name_col = c
                    break
            if pts_col and name_col:
                df = df.fillna({pts_col: 0})
                df_sorted = df.sort_values(by=pts_col, ascending=False)
                rows = []
                for _, r in df_sorted.head(10).iterrows():
                    rows.append((r.get(name_col, "Unknown"), int(r.get(pts_col, 0))))
                chart_png = output.parent / "top_points_supabase_inferred.png"
                make_points_chart(rows, chart_png)
                add_image_slide(prs, "Top Cadets — Points (inferred)", chart_png)
                add_table_slide(prs, "Top Cadets Table (inferred)", rows, col_names=["Name", "Points"])

    prs.save(str(output))
    return output


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", default="outputs/cadets_presentation_supabase.pptx")
    parser.add_argument("--project-id", default=None)
    parser.add_argument("--function-name", default="server")
    args = parser.parse_args()

    project_id = args.project_id or os.environ.get("SUPABASE_PROJECT_ID") or read_project_id_from_repo()
    if not project_id:
        print("Error: SUPABASE project id not found. Set SUPABASE_PROJECT_ID env var or add utils/supabase/info.tsx")
        return

    api_key = os.environ.get("SUPABASE_SERVICE_ROLE_KEY") or os.environ.get("SUPABASE_ANON_KEY")

    out = Path(args.output)
    result = build_presentation_from_supabase(out, project_id, args.function_name, api_key)
    print(f"Generated: {result}")


if __name__ == "__main__":
    main()
